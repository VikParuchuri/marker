import time
import torch
import os
from datetime import datetime
from marker.utils import send_callback, flush_cuda_memory
import pytz
# 获取北京时区
beijing_tz = pytz.timezone('Asia/Shanghai')
os.environ["PYTORCH_ENABLE_MPS_FALLBACK"] = "1" # For some reason, transformers decided to use .isin for a simple op, which is not supported on MPS

# from marker.convert import convert_single_pdf
from marker.config.parser import ConfigParser
from marker.converters.pdf import PdfConverter
from marker.settings import settings
from marker.logger import configure_logging
from marker.models import create_model_dict

configure_logging()


class ExtractionProc:
    def __init__(self, config=None):
        self.model_lst = []
    
    def load_models(self):
        # self.model_lst = load_all_models(torch.device("cuda"), dtype=torch.float16)
        self.model_dict = create_model_dict()
    
    def parse_docx(self, file_byte):
        # 使用python-docx库解析docx文件
        from io import BytesIO
        from docx import Document
        
        # 创建BytesIO对象
        docx_file = BytesIO(file_byte)
        
        # 加载docx文件
        doc = Document(docx_file)
        
        # 提取所有段落中的文本
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
                
        # 提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text)
        
        # 将所有文本拼接成一个字符串
        text = "\n".join(full_text)
        return text, None, None

    def parse_pptx(self, file_byte):
        # 使用python-pptx库解析pptx文件
        from io import BytesIO
        from pptx import Presentation
        
        # 创建BytesIO对象
        pptx_file = BytesIO(file_byte)
        
        # 加载pptx文件
        prs = Presentation(pptx_file)
        
        # 提取所有幻灯片中的文本
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        
        # 将所有文本拼接成一个字符串
        text = "\n".join(full_text)
        return text, None, None

    def parse_txt(self, file_byte):
        # 这个不需解析，返回文本
        full_text = file_byte.decode('utf-8')
        return full_text, None, None

    def extraction(self, args, file_byte, callback_url='', docId='', file_type='pdf'):
        start = time.time()
        kwargs = {
            'output_dir': args.get('output_dir', settings.OUTPUT_DIR),
            'debug': args.get('debug', False),
            'output_format': args.get('output_format', 'markdown'),
            'page_range': args.get('page_range', None),
            'force_ocr': args.get('force_ocr', False),
            'processors': args.get('processors', None),
            'config_json': args.get('config_json', None),
            'languages': args.get('languages', None),
            'disable_multiprocessing': args.get('disable_multiprocessing', False),
            'paginate_output': args.get('paginate_output', False),
            'disable_image_extraction': args.get('disable_image_extraction', False),
        }
        config_parser = ConfigParser(kwargs)
        time_str = datetime.now(beijing_tz).strftime("%H:%M:%S")
        send_callback(callback_url, {
            'status': True,
            'messages': 'success',
            'docId': docId,
            'progress': 5,
            'progress_text': '开始解析  ' + time_str
        })
        converter = PdfConverter(
            config=config_parser.generate_config_dict(),
            artifact_dict=self.model_dict,
            processor_list=config_parser.get_processors(),
            renderer=config_parser.get_renderer(),
            callback_url=callback_url,
            docId=docId,
            llm_service="marker.services.gemini.GoogleGeminiService"
        )
        rendered = converter(file_byte, file_type=file_type)
        # 统计metadata中有多少page 多少表格 多少公式和ocr次数
        info = self.count_table_formula(rendered.metadata)
        full_text = rendered.markdown
        print('use time: ', time.time() - start)

        return full_text, None, info, rendered.metadata
    
    def count_table_formula(self, metadata):
        table_count = 0
        formula_count = 0
        ocr_count = 0
        for page in metadata['page_stats']:
            meta = page['block_counts']
            extract_method = page['text_extraction_method']
            for i in meta:
                if i[0] == 'Table':
                    table_count += i[1]
                elif i[0] == 'Equation':
                    formula_count += i[1]
            if extract_method != 'pdftext':
                ocr_count += 1
        return {
            'table_count': table_count,
            'formula_count': formula_count,
            'ocr_count': ocr_count
        }
    
    def split_chunk(self):
        pass

    def annotate_chunk(self):
        pass
