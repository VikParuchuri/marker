import base64
import os
import tempfile
import traceback

from marker.logger import get_logger
from marker.providers.pdf import PdfProvider

logger = get_logger()

css = """
@page {
    size: A4 landscape;
    margin: 1.5cm;
}

table {
    width: 100%;
    border-collapse: collapse;
    break-inside: auto;
    font-size: 10pt;
}

tr {
    break-inside: avoid;
    page-break-inside: avoid;
}

td {
    border: 0.75pt solid #000;
    padding: 6pt;
}

img {
    max-width: 100%;
    height: auto;
    object-fit: contain;
}
"""


class PowerPointProvider(PdfProvider):
    include_slide_number: bool = False

    def __init__(self, filepath: str, config=None):
        temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        self.temp_pdf_path = temp_pdf.name
        temp_pdf.close()

        # Convert PPTX to PDF
        try:
            self.convert_pptx_to_pdf(filepath)
        except Exception as e:
            print(traceback.format_exc())
            raise ValueError(f"Error converting PPTX to PDF: {e}")

        # Initalize the PDF provider with the temp pdf path
        super().__init__(self.temp_pdf_path, config)

    def __del__(self):
        if os.path.exists(self.temp_pdf_path):
            os.remove(self.temp_pdf_path)

    def convert_pptx_to_pdf(self, filepath):
        from weasyprint import CSS, HTML
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        pptx = Presentation(filepath)

        html_parts = []

        for slide_index, slide in enumerate(pptx.slides):
            html_parts.append("<section>")
            if self.include_slide_number:
                html_parts.append(f"<h2>Slide {slide_index + 1}</h2>")

            # Process shapes in the slide
            for shape in slide.shapes:
                # If shape is a group shape, we recursively handle all grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    html_parts.append(self._handle_group(shape))
                    continue

                # If shape is a table
                if shape.has_table:
                    html_parts.append(self._handle_table(shape))
                    continue

                # If shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    html_parts.append(self._handle_image(shape))
                    continue

                # If shape has text
                if hasattr(shape, "text") and shape.text is not None:
                    if shape.has_text_frame:
                        # Distinguish placeholders (title, subtitle, etc.)
                        html_parts.append(self._handle_text(shape))
                    else:
                        html_parts.append(f"<p>{self._escape_html(shape.text)}</p>")

            html_parts.append("</section>")

        html = "\n".join(html_parts)

        # We convert the HTML into a PDF
        HTML(string=html).write_pdf(
            self.temp_pdf_path, stylesheets=[CSS(string=css), self.get_font_css()]
        )

    def _handle_group(self, group_shape) -> str:
        """
        Recursively handle shapes in a group. Returns HTML string for the entire group.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        group_parts = []
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_parts.append(self._handle_group(shape))
                continue

            if shape.has_table:
                group_parts.append(self._handle_table(shape))
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                group_parts.append(self._handle_image(shape))
                continue

            if hasattr(shape, "text"):
                if shape.has_text_frame:
                    group_parts.append(self._handle_text(shape))
                else:
                    group_parts.append(f"<p>{self._escape_html(shape.text)}</p>")

        return "".join(group_parts)

    def _handle_text(self, shape) -> str:
        """
        Processes shape text, including bullet/numbered list detection and placeholders
        (title, subtitle, etc.). Returns HTML for the text block(s).
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        # Distinguish placeholders to see if it's a title or subtitle
        label_html_tag = "p"
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                label_html_tag = "h3"
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                label_html_tag = "h4"

        # Keep track of whether we are currently in a <ul> or <ol>
        html_parts = []
        list_open = False
        list_type = None  # "ul" or "ol"

        for paragraph in shape.text_frame.paragraphs:
            p_el = paragraph._element
            # Check bullet
            bullet_char = p_el.find(".//a:buChar", namespaces=p_el.nsmap)
            bullet_num = p_el.find(".//a:buAutoNum", namespaces=p_el.nsmap)

            is_bullet = (bullet_char is not None) or (paragraph.level > 0)
            is_numbered = bullet_num is not None

            # If the paragraph is bullet or numbered
            if is_bullet or is_numbered:
                # Decide if we need to start a new list or continue an existing one
                current_list_type = "ol" if is_numbered else "ul"
                if not list_open:
                    # Start new
                    list_open = True
                    list_type = current_list_type
                    html_parts.append(f"<{list_type}>")

                elif list_open and list_type != current_list_type:
                    # Close old list, start new
                    html_parts.append(f"</{list_type}>")
                    list_type = current_list_type
                    html_parts.append(f"<{list_type}>")

                # Build the bullet (li) text from all runs in the paragraph
                p_text = "".join(run.text for run in paragraph.runs)
                if p_text:
                    html_parts.append(f"<li>{self._escape_html(p_text)}</li>")

            else:
                # If we were in a list, we need to close it
                if list_open:
                    html_parts.append(f"</{list_type}>")
                    list_open = False
                    list_type = None

                # Now it's just a normal paragraph
                # Gather the paragraph text from runs
                p_text = "".join(run.text for run in paragraph.runs)
                if p_text:
                    # If we know it's a slide title, we can use <h3> or so
                    html_parts.append(
                        f"<{label_html_tag}>{self._escape_html(p_text)}</{label_html_tag}>"
                    )

        # If the text frame ended and we still have an open list, close it
        if list_open:
            html_parts.append(f"</{list_type}>")

        return "".join(html_parts)

    def _handle_image(self, shape) -> str:
        """
        Embeds the image as a base64 <img> in HTML.
        """
        image = shape.image
        image_bytes = image.blob

        try:
            img_str = base64.b64encode(image_bytes).decode("utf-8")
            return f"<img src='data:{image.content_type};base64,{img_str}' />"
        except Exception as e:
            logger.warning(f"Warning: image cannot be loaded by Pillow: {e}")
            return ""

    def _handle_table(self, shape) -> str:
        """
        Renders a shape's table as an HTML <table>.
        """
        table_html = []
        table_html.append("<table border='1'>")

        for row in shape.table.rows:
            row_html = ["<tr>"]
            for cell in row.cells:
                row_html.append(f"<td>{self._escape_html(cell.text)}</td>")
            row_html.append("</tr>")
            table_html.append("".join(row_html))

        table_html.append("</table>")
        return "".join(table_html)

    def _escape_html(self, text: str) -> str:
        """
        Minimal escaping for HTML special characters.
        """
        return (
            text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;")
        )
